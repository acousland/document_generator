"""PowerPoint document generator using python-pptx."""

from pathlib import Path
from typing import Dict, Any
import re
from pptx import Presentation
from .base import DocumentGenerator


class PowerPointGenerator(DocumentGenerator):
    """Generator for Microsoft PowerPoint documents."""
    
    def generate(self, template_path: Path, fields: Dict[str, Any], output_path: Path) -> Path:
        """
        Generate a PowerPoint document from a template by replacing placeholders.
        
        Placeholders in the template should be in the format {{field_name}}.
        """
        # Load the template
        prs = Presentation(template_path)
        
        # Iterate through all slides
        for slide in prs.slides:
            # Replace in shapes
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    self._replace_text_in_shape(shape, fields)
                # Replace in tables
                if hasattr(shape, "table"):
                    table = shape.table
                    for row in table.rows:
                        for cell in row.cells:
                            self._replace_text_in_shape(cell, fields)
        
        # Save the presentation
        prs.save(output_path)
        return output_path
    
    def _replace_text_in_shape(self, shape, fields: Dict[str, Any]) -> None:
        """Replace placeholder text in a shape."""
        if not hasattr(shape, "text_frame"):
            return
            
        text_frame = shape.text_frame
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                for field_name, field_value in fields.items():
                    placeholder = f"{{{{{field_name}}}}}"
                    if placeholder in run.text:
                        run.text = run.text.replace(placeholder, str(field_value))
    
    def get_template_fields(self, template_path: Path) -> list[str]:
        """Extract field names from a PowerPoint template."""
        prs = Presentation(template_path)
        fields = set()
        
        # Find placeholders in format {{field_name}}
        pattern = r'\{\{([^}]+)\}\}'
        
        # Search in all slides
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    text = shape.text
                    matches = re.findall(pattern, text)
                    fields.update(matches)
                # Search in tables
                if hasattr(shape, "table"):
                    table = shape.table
                    for row in table.rows:
                        for cell in row.cells:
                            if hasattr(cell, "text_frame"):
                                text = cell.text
                                matches = re.findall(pattern, text)
                                fields.update(matches)
        
        return sorted(list(fields))
