"""Script to create example templates for testing."""

from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches, Pt
from pathlib import Path


def create_word_template():
    """Create an example Word template."""
    doc = Document()
    
    # Add title
    doc.add_heading('Letter Template', 0)
    
    # Add content with placeholders
    doc.add_paragraph(f'Date: {{{{date}}}}')
    doc.add_paragraph()
    
    doc.add_paragraph(f'Dear {{{{recipient_name}}}},')
    doc.add_paragraph()
    
    doc.add_paragraph(f'{{{{body_text}}}}')
    doc.add_paragraph()
    
    doc.add_paragraph('Sincerely,')
    doc.add_paragraph(f'{{{{sender_name}}}}')
    doc.add_paragraph(f'{{{{sender_title}}}}')
    
    # Add a table
    doc.add_paragraph()
    table = doc.add_table(rows=3, cols=2)
    table.style = 'Light Grid Accent 1'
    
    cells = table.rows[0].cells
    cells[0].text = 'Field'
    cells[1].text = 'Value'
    
    cells = table.rows[1].cells
    cells[0].text = 'Company'
    cells[1].text = '{{company_name}}'
    
    cells = table.rows[2].cells
    cells[0].text = 'Phone'
    cells[1].text = '{{phone_number}}'
    
    # Save template
    template_path = Path('templates/letter.docx')
    template_path.parent.mkdir(parents=True, exist_ok=True)
    doc.save(template_path)
    print(f"Created Word template: {template_path}")


def create_excel_template():
    """Create an example Excel template."""
    wb = Workbook()
    ws = wb.active
    ws.title = "Report"
    
    # Add headers
    ws['A1'] = 'Monthly Report'
    ws['A1'].font = ws['A1'].font.copy(bold=True, size=14)
    
    ws['A3'] = 'Report Date:'
    ws['B3'] = '{{report_date}}'
    
    ws['A4'] = 'Department:'
    ws['B4'] = '{{department}}'
    
    ws['A6'] = 'Metric'
    ws['B6'] = 'Value'
    ws['A6'].font = ws['A6'].font.copy(bold=True)
    ws['B6'].font = ws['B6'].font.copy(bold=True)
    
    ws['A7'] = 'Total Sales'
    ws['B7'] = '{{total_sales}}'
    
    ws['A8'] = 'New Customers'
    ws['B8'] = '{{new_customers}}'
    
    ws['A9'] = 'Revenue'
    ws['B9'] = '{{revenue}}'
    
    ws['A11'] = 'Notes:'
    ws['A12'] = '{{notes}}'
    
    # Adjust column widths
    ws.column_dimensions['A'].width = 20
    ws.column_dimensions['B'].width = 30
    
    # Save template
    template_path = Path('templates/report.xlsx')
    template_path.parent.mkdir(parents=True, exist_ok=True)
    wb.save(template_path)
    print(f"Created Excel template: {template_path}")


def create_powerpoint_template():
    """Create an example PowerPoint template."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "{{presentation_title}}"
    subtitle.text = "{{presenter_name}}\n{{date}}"
    
    # Slide 2: Content Slide
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "{{slide_title}}"
    
    text_frame = body_shape.text_frame
    text_frame.text = "{{bullet_1}}"
    
    p = text_frame.add_paragraph()
    p.text = "{{bullet_2}}"
    p.level = 0
    
    p = text_frame.add_paragraph()
    p.text = "{{bullet_3}}"
    p.level = 0
    
    # Slide 3: Data Slide
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    
    left = Inches(1)
    top = Inches(1)
    width = Inches(8)
    height = Inches(1)
    
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "Summary: {{summary}}"
    
    # Save template
    template_path = Path('templates/presentation.pptx')
    template_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(template_path)
    print(f"Created PowerPoint template: {template_path}")


if __name__ == "__main__":
    create_word_template()
    create_excel_template()
    create_powerpoint_template()
    print("\nAll templates created successfully!")
