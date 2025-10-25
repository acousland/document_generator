"""Create an example PowerPoint template with metadata in slide notes."""

import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pathlib import Path


def create_advanced_presentation_template():
    """Create a PowerPoint template with metadata-driven slide layouts."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title Page
    blank_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(blank_layout)
    
    # Add title
    left = Inches(1)
    top = Inches(2.5)
    width = Inches(8)
    height = Inches(1)
    
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = "{{title}}"
    title_frame.paragraphs[0].font.size = Pt(44)
    title_frame.paragraphs[0].font.bold = True
    
    # Add subtitle
    subtitle_box = slide.shapes.add_textbox(left, Inches(3.5), width, Inches(0.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "{{subtitle}}"
    subtitle_frame.paragraphs[0].font.size = Pt(24)
    
    # Add presenter
    presenter_box = slide.shapes.add_textbox(left, Inches(4.5), width, Inches(0.5))
    presenter_frame = presenter_box.text_frame
    presenter_frame.text = "{{presenter}}"
    presenter_frame.paragraphs[0].font.size = Pt(18)
    
    # Add metadata to notes
    notes_slide = slide.notes_slide
    notes_text_frame = notes_slide.notes_text_frame
    metadata = {
        "slide_type": "title_page",
        "description": "Main title slide for presentation opening",
        "placeholders": {
            "title": {
                "type": "text",
                "description": "Main presentation title"
            },
            "subtitle": {
                "type": "text",
                "description": "Presentation subtitle or date"
            },
            "presenter": {
                "type": "text",
                "description": "Presenter name"
            }
        }
    }
    notes_text_frame.text = json.dumps(metadata, indent=2)
    
    # Slide 2: Section Break
    slide = prs.slides.add_slide(blank_layout)
    
    section_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1.5))
    section_frame = section_box.text_frame
    section_frame.text = "{{section_title}}"
    section_frame.paragraphs[0].font.size = Pt(54)
    section_frame.paragraphs[0].font.bold = True
    
    notes_slide = slide.notes_slide
    notes_text_frame = notes_slide.notes_text_frame
    metadata = {
        "slide_type": "section_break",
        "description": "Section divider slide",
        "placeholders": {
            "section_title": {
                "type": "text",
                "description": "Section heading"
            }
        }
    }
    notes_text_frame.text = json.dumps(metadata, indent=2)
    
    # Slide 3: Content Slide
    slide = prs.slides.add_slide(blank_layout)
    
    # Heading
    heading_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(0.75))
    heading_frame = heading_box.text_frame
    heading_frame.text = "{{heading}}"
    heading_frame.paragraphs[0].font.size = Pt(36)
    heading_frame.paragraphs[0].font.bold = True
    
    # Body text
    body_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(2))
    body_frame = body_box.text_frame
    body_frame.text = "{{body}}"
    body_frame.paragraphs[0].font.size = Pt(18)
    body_frame.word_wrap = True
    
    # Bullet points
    bullets_box = slide.shapes.add_textbox(Inches(1.5), Inches(4), Inches(7), Inches(2.5))
    bullets_frame = bullets_box.text_frame
    bullets_frame.text = "{{bullet_points}}"
    bullets_frame.paragraphs[0].font.size = Pt(16)
    
    notes_slide = slide.notes_slide
    notes_text_frame = notes_slide.notes_text_frame
    metadata = {
        "slide_type": "content",
        "description": "Standard content slide with heading and body",
        "placeholders": {
            "heading": {
                "type": "text",
                "description": "Slide heading"
            },
            "body": {
                "type": "paragraph",
                "description": "Main content text"
            },
            "bullet_points": {
                "type": "list",
                "description": "List of key points (array of strings)"
            }
        }
    }
    notes_text_frame.text = json.dumps(metadata, indent=2)
    
    # Slide 4: Two Column Layout
    slide = prs.slides.add_slide(blank_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(0.75))
    title_frame = title_box.text_frame
    title_frame.text = "{{heading}}"
    title_frame.paragraphs[0].font.size = Pt(36)
    title_frame.paragraphs[0].font.bold = True
    
    # Left column
    left_heading_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4), Inches(0.5))
    left_heading_frame = left_heading_box.text_frame
    left_heading_frame.text = "{{left_heading}}"
    left_heading_frame.paragraphs[0].font.size = Pt(24)
    left_heading_frame.paragraphs[0].font.bold = True
    
    left_body_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(4), Inches(4))
    left_body_frame = left_body_box.text_frame
    left_body_frame.text = "{{left_content}}"
    left_body_frame.paragraphs[0].font.size = Pt(16)
    left_body_frame.word_wrap = True
    
    # Right column
    right_heading_box = slide.shapes.add_textbox(Inches(5.5), Inches(1.5), Inches(4), Inches(0.5))
    right_heading_frame = right_heading_box.text_frame
    right_heading_frame.text = "{{right_heading}}"
    right_heading_frame.paragraphs[0].font.size = Pt(24)
    right_heading_frame.paragraphs[0].font.bold = True
    
    right_body_box = slide.shapes.add_textbox(Inches(5.5), Inches(2.5), Inches(4), Inches(4))
    right_body_frame = right_body_box.text_frame
    right_body_frame.text = "{{right_content}}"
    right_body_frame.paragraphs[0].font.size = Pt(16)
    right_body_frame.word_wrap = True
    
    notes_slide = slide.notes_slide
    notes_text_frame = notes_slide.notes_text_frame
    metadata = {
        "slide_type": "two_column",
        "description": "Two column layout for comparisons or parallel content",
        "placeholders": {
            "heading": {
                "type": "text",
                "description": "Slide heading"
            },
            "left_heading": {
                "type": "text",
                "description": "Left column heading"
            },
            "left_content": {
                "type": "paragraph",
                "description": "Left column content"
            },
            "right_heading": {
                "type": "text",
                "description": "Right column heading"
            },
            "right_content": {
                "type": "paragraph",
                "description": "Right column content"
            }
        }
    }
    notes_text_frame.text = json.dumps(metadata, indent=2)
    
    # Slide 5: Closing Slide
    slide = prs.slides.add_slide(blank_layout)
    
    closing_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
    closing_frame = closing_box.text_frame
    closing_frame.text = "{{closing_message}}"
    closing_frame.paragraphs[0].font.size = Pt(36)
    
    contact_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(1))
    contact_frame = contact_box.text_frame
    contact_frame.text = "{{contact_info}}"
    contact_frame.paragraphs[0].font.size = Pt(18)
    
    notes_slide = slide.notes_slide
    notes_text_frame = notes_slide.notes_text_frame
    metadata = {
        "slide_type": "closing",
        "description": "Closing slide with message and contact information",
        "placeholders": {
            "closing_message": {
                "type": "text",
                "description": "Closing message or call to action"
            },
            "contact_info": {
                "type": "text",
                "description": "Contact information"
            }
        }
    }
    notes_text_frame.text = json.dumps(metadata, indent=2)
    
    # Save template
    template_path = Path('templates/advanced_presentation.pptx')
    template_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(template_path)
    print(f"✅ Created advanced PowerPoint template: {template_path}")
    print("\nSlide types included:")
    print("  - title_page: Title, subtitle, presenter")
    print("  - section_break: Section heading")
    print("  - content: Heading, body, bullet points")
    print("  - two_column: Side-by-side content")
    print("  - closing: Closing message and contact info")


if __name__ == "__main__":
    create_advanced_presentation_template()
