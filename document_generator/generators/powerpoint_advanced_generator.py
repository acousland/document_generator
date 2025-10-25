"""Advanced PowerPoint generator with slide composition capabilities."""

from pathlib import Path
from typing import Dict, Any, List, Optional
import re
import json
from pptx import Presentation
from pptx.util import Pt
from copy import deepcopy
from .base import DocumentGenerator


class PowerPointAdvancedGenerator(DocumentGenerator):
    """Advanced generator for composing PowerPoint presentations from slide layouts with metadata."""
    
    def parse_slide_metadata(self, slide) -> Optional[Dict[str, Any]]:
        """
        Parse JSON metadata from slide notes.
        
        Returns None if no valid metadata found.
        """
        if not slide.has_notes_slide:
            return None
        
        notes_text = slide.notes_slide.notes_text_frame.text.strip()
        if not notes_text:
            return None
        
        # Replace smart/curly quotes with straight quotes
        # PowerPoint often auto-converts straight quotes to smart quotes
        notes_text = notes_text.replace('\u201c', '"').replace('\u201d', '"')  # Left (") and right (") double quotes
        notes_text = notes_text.replace('\u2018', "'").replace('\u2019', "'")  # Left (') and right (') single quotes
        
        try:
            metadata = json.loads(notes_text)
            return metadata
        except json.JSONDecodeError:
            return None
    
    def get_template_slide_types(self, template_path: Path) -> List[Dict[str, Any]]:
        """
        Get all slide types defined in the template with their metadata.
        
        Returns:
            List of slide type information including slide_type, description, 
            placeholders, and slide_index.
        """
        prs = Presentation(template_path)
        slide_types = []
        
        for idx, slide in enumerate(prs.slides):
            metadata = self.parse_slide_metadata(slide)
            if metadata and "slide_type" in metadata:
                slide_info = {
                    "slide_index": idx,
                    "slide_type": metadata.get("slide_type"),
                    "description": metadata.get("description", ""),
                    "placeholders": metadata.get("placeholders", {}),
                }
                slide_types.append(slide_info)
        
        return slide_types
    
    def generate_from_slides(
        self, 
        template_path: Path, 
        slides_data: List[Dict[str, Any]], 
        output_path: Path
    ) -> Path:
        """
        Generate a PowerPoint presentation by composing slides based on slide_type.
        
        Args:
            template_path: Path to the template with metadata in slide notes
            slides_data: List of slide specifications, each containing:
                - slide_type: Type of slide to use (matches metadata)
                - fields: Dictionary of field names and values for that slide
            output_path: Where to save the generated presentation
            
        Example slides_data:
            [
                {
                    "slide_type": "title_page",
                    "fields": {"title": "My Presentation", "subtitle": "Oct 2025"}
                },
                {
                    "slide_type": "content",
                    "fields": {"heading": "Section 1", "body": "Content here"}
                }
            ]
        """
        # Load the template
        template_prs = Presentation(template_path)
        
        # Build a map of slide_type -> list of slide objects with that type
        slide_type_map = {}
        for idx, slide in enumerate(template_prs.slides):
            metadata = self.parse_slide_metadata(slide)
            if metadata and "slide_type" in metadata:
                slide_type = metadata["slide_type"]
                if slide_type not in slide_type_map:
                    slide_type_map[slide_type] = []
                slide_type_map[slide_type].append({
                    'index': idx,
                    'slide': slide,
                    'metadata': metadata
                })
        
        # Build list of slides to keep (in order) with their data
        slides_to_build = []
        for slide_spec in slides_data:
            slide_type = slide_spec.get("slide_type")
            fields = slide_spec.get("fields", {})
            
            # Find the template slide with this type
            if slide_type not in slide_type_map:
                raise ValueError(
                    f"slide_type '{slide_type}' not found in template. "
                    f"Available types: {list(slide_type_map.keys())}"
                )
            
            # If there are multiple slides with this type, pick the one that best matches the fields
            candidate_slides = slide_type_map[slide_type]
            
            if len(candidate_slides) == 1:
                # Only one option
                template_info = candidate_slides[0]
            else:
                # Multiple slides with same type - pick based on field match
                best_match = None
                best_score = -1
                
                for candidate in candidate_slides:
                    candidate_fields = set(candidate['metadata'].get('placeholders', {}).keys())
                    requested_fields = set(fields.keys())
                    
                    # Count matching fields
                    matches = len(candidate_fields & requested_fields)
                    
                    # Prefer candidates that have most of the requested fields
                    if matches > best_score:
                        best_score = matches
                        best_match = candidate
                
                template_info = best_match if best_match else candidate_slides[0]
            
            slides_to_build.append({
                'template_index': template_info['index'],
                'fields': fields,
                'metadata': template_info['metadata']
            })
        
        # Create new presentation from template (preserves themes, layouts, etc.)
        import shutil
        shutil.copy(template_path, output_path)
        output_prs = Presentation(output_path)
        
        # Delete all existing slides (in reverse)
        slide_count = len(output_prs.slides)
        for i in range(slide_count - 1, -1, -1):
            rId = output_prs.slides._sldIdLst[i].rId
            output_prs.part.drop_rel(rId)
            del output_prs.slides._sldIdLst[i]
        
        # Reload template to get fresh slide references
        template_prs = Presentation(template_path)
        
        # Add slides by copying content from template slides
        for build_info in slides_to_build:
            template_idx = build_info['template_index']
            template_slide = template_prs.slides[template_idx]
            
            # Find the layout by name in the output presentation
            layout_name = template_slide.slide_layout.name
            matching_layout = None
            for layout in output_prs.slide_layouts:
                if layout.name == layout_name:
                    matching_layout = layout
                    break
            
            if matching_layout is None:
                # Fallback: try to use layout at same index, or blank layout
                try:
                    layout_idx = template_slide.slide_layout.element.getparent().index(template_slide.slide_layout.element)
                    if layout_idx < len(output_prs.slide_layouts):
                        matching_layout = output_prs.slide_layouts[layout_idx]
                except:
                    pass
                
                if matching_layout is None:
                    # Last resort: use first layout or blank
                    matching_layout = output_prs.slide_layouts[0] if len(output_prs.slide_layouts) > 0 else None
            
            if matching_layout is None:
                raise ValueError(f"Could not find matching layout for '{layout_name}'")
            
            # Add slide using the layout
            new_slide = output_prs.slides.add_slide(matching_layout)
            
            # Copy all shapes from the template slide (images, custom shapes, etc.)
            self._copy_slide_shapes(template_slide, new_slide)
            
            # Copy placeholder content from template (title, body, etc.)
            self._copy_placeholder_content(template_slide, new_slide)
            
            # Now populate with data (this will replace placeholders)
            self._populate_slide(new_slide, build_info['fields'], build_info['metadata'])
        
        # Save the presentation
        output_prs.save(output_path)
        return output_path
    
    def _copy_placeholder_content(self, source_slide, target_slide):
        """
        Copy placeholder text content from source slide to target slide.
        This ensures placeholders have their template content before we replace them.
        """
        for source_shape in source_slide.shapes:
            # Only copy from placeholders with text
            if hasattr(source_shape, 'is_placeholder') and source_shape.is_placeholder:
                if not hasattr(source_shape, 'text_frame'):
                    continue
                
                # Find the corresponding placeholder in target by index
                try:
                    source_idx = source_shape.placeholder_format.idx
                    for target_shape in target_slide.placeholders:
                        if target_shape.placeholder_format.idx == source_idx:
                            # Copy text content from source to target
                            if hasattr(target_shape, 'text_frame') and hasattr(source_shape, 'text_frame'):
                                self._copy_text_frame_content(source_shape.text_frame, target_shape.text_frame)
                            break
                except:
                    pass
    
    def _copy_slide_shapes(self, source_slide, target_slide):
        """
        Copy all shapes from source slide to target slide.
        This includes images, text boxes, and other custom content.
        Fills picture placeholders if they contain images.
        Skips text placeholders since they come from the layout.
        """
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        from pptx.util import Inches
        import io
        
        # First, identify picture placeholders in source that have images
        # and fill corresponding placeholders in target
        for source_shape in source_slide.shapes:
            is_placeholder = hasattr(source_shape, 'is_placeholder') and source_shape.is_placeholder
            
            if is_placeholder and source_shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                # Check if it's a picture placeholder with an image
                try:
                    if hasattr(source_shape, 'image'):
                        source_image = source_shape.image
                        
                        # Find the corresponding placeholder in target by index
                        source_idx = source_shape.placeholder_format.idx
                        target_placeholder = None
                        for target_shape in target_slide.placeholders:
                            if target_shape.placeholder_format.idx == source_idx:
                                target_placeholder = target_shape
                                break
                        
                        if target_placeholder is not None:
                            # Insert picture into the placeholder
                            try:
                                target_placeholder.insert_picture(io.BytesIO(source_image.blob))
                            except:
                                # If insert fails, add as separate picture
                                target_slide.shapes.add_picture(
                                    io.BytesIO(source_image.blob),
                                    source_shape.left,
                                    source_shape.top,
                                    source_shape.width,
                                    source_shape.height
                                )
                except:
                    pass
        
        # Now copy non-placeholder shapes
        for shape in source_slide.shapes:
            # Skip all placeholders - we handled picture placeholders above
            if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
                continue
            
            # Skip decorative auto shapes (rectangles, lines, etc.) that don't have text
            # These are often covering lines or used for spacing and theme colors don't transfer well
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                # Skip unless it has text content
                if not (hasattr(shape, 'text_frame') and shape.text_frame.text.strip()):
                    continue
            
            # Skip connector lines and other connector shapes (shape_type 9 = LINE, 10 = CONNECTOR)
            if shape.shape_type in [9, 10]:
                continue
            
            # Copy based on shape type
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                # Copy picture
                try:
                    # Get the image from the source shape
                    image = shape.image
                    image_bytes = image.blob
                    
                    # Add the image to the target slide at the same position
                    picture = target_slide.shapes.add_picture(
                        io.BytesIO(image_bytes),
                        shape.left,
                        shape.top,
                        shape.width,
                        shape.height
                    )
                except Exception as e:
                    # If image copy fails, skip it
                    pass
                    
            elif shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                # Copy text box
                try:
                    text_box = target_slide.shapes.add_textbox(
                        shape.left,
                        shape.top,
                        shape.width,
                        shape.height
                    )
                    # Copy text and formatting
                    if hasattr(shape, 'text_frame'):
                        self._copy_text_frame_content(shape.text_frame, text_box.text_frame)
                except:
                    pass
                    
            elif shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE or shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                # Copy auto shapes and grouped shapes
                try:
                    # For auto shapes, try to duplicate with same properties
                    if hasattr(shape, 'auto_shape_type'):
                        new_shape = target_slide.shapes.add_shape(
                            shape.auto_shape_type,
                            shape.left,
                            shape.top,
                            shape.width,
                            shape.height
                        )
                        # Copy fill and line properties if possible
                        if hasattr(shape, 'fill') and hasattr(new_shape, 'fill'):
                            try:
                                if shape.fill.type == 1:  # Solid fill
                                    new_shape.fill.solid()
                                    new_shape.fill.fore_color.rgb = shape.fill.fore_color.rgb
                            except:
                                pass
                        # Copy text if it has text
                        if hasattr(shape, 'text_frame') and hasattr(new_shape, 'text_frame'):
                            self._copy_text_frame_content(shape.text_frame, new_shape.text_frame)
                except:
                    pass
    
    def _copy_text_frame_content(self, source_frame, target_frame):
        """Copy text content and basic formatting from source to target text frame."""
        try:
            target_frame.clear()
            for source_para in source_frame.paragraphs:
                target_para = target_frame.add_paragraph()
                target_para.text = source_para.text
                target_para.level = source_para.level
                
                # Copy alignment if set
                if source_para.alignment is not None:
                    target_para.alignment = source_para.alignment
        except:
            # If copying fails, just set the plain text
            try:
                target_frame.text = source_frame.text
            except:
                pass
    
    def _populate_slide(self, slide, fields: Dict[str, Any], metadata: Optional[Dict[str, Any]]):
        """
        Populate slide with field values based on their types defined in metadata.
        """
        placeholders_info = metadata.get("placeholders", {}) if metadata else {}
        
        # First pass: handle simple text replacements
        for field_name, field_value in fields.items():
            placeholder_info = placeholders_info.get(field_name, {})
            field_type = placeholder_info.get("type", "text")
            
            if field_type in ["text", "paragraph", "number", "date"]:
                # Simple replacement
                placeholder = f"{{{{{field_name}}}}}"
                self._replace_text_in_slide(slide, placeholder, str(field_value))
            
            elif field_type == "list":
                # Handle bullet points
                self._populate_list_field(slide, field_name, field_value)
            
            elif field_type == "table":
                # Handle table data
                self._populate_table_field(slide, field_name, field_value)
            
            elif field_type == "image":
                # Handle image insertion
                self._populate_image_field(slide, field_name, field_value)
    
    def _replace_text_in_slide(self, slide, placeholder: str, value: str):
        """Replace placeholder text in all text elements of a slide."""
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                # First, try direct replacement
                text_frame = shape.text_frame
                
                # Collect all runs in order
                all_text = shape.text
                if placeholder in all_text:
                    # Replace in the text frame
                    # Since PowerPoint might split the text across runs,
                    # we need to rebuild the text
                    new_text = all_text.replace(placeholder, value)
                    
                    # Clear and rebuild
                    text_frame.clear()
                    p = text_frame.paragraphs[0]
                    p.text = new_text
            
            # Handle tables
            if hasattr(shape, "table"):
                table = shape.table
                for row in table.rows:
                    for cell in row.cells:
                        if hasattr(cell, "text_frame"):
                            cell_text = cell.text
                            if placeholder in cell_text:
                                new_text = cell_text.replace(placeholder, value)
                                cell.text_frame.clear()
                                p = cell.text_frame.paragraphs[0]
                                p.text = new_text
    
    def _populate_list_field(self, slide, field_name: str, items: List[str]):
        """
        Populate a bullet point list field.
        Finds a shape containing {{field_name}} and replaces it with bullet points.
        """
        placeholder = f"{{{{{field_name}}}}}"
        
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                # Check if this shape contains the placeholder
                if placeholder in shape.text:
                    text_frame = shape.text_frame
                    text_frame.clear()
                    
                    # Add each item as a bullet point
                    for i, item in enumerate(items):
                        if i == 0:
                            p = text_frame.paragraphs[0]
                        else:
                            p = text_frame.add_paragraph()
                        
                        p.text = str(item)
                        p.level = 0
                    
                    return
    
    def _populate_table_field(self, slide, field_name: str, table_data: List[List[Any]]):
        """
        Populate a table field with data.
        Expects table_data as a 2D array: [[row1_col1, row1_col2], [row2_col1, row2_col2], ...]
        """
        placeholder = f"{{{{{field_name}}}}}"
        
        for shape in slide.shapes:
            if hasattr(shape, "table"):
                table = shape.table
                
                # Check if table contains the placeholder
                has_placeholder = False
                for row in table.rows:
                    for cell in row.cells:
                        if placeholder in cell.text:
                            has_placeholder = True
                            break
                    if has_placeholder:
                        break
                
                if has_placeholder and table_data:
                    # Clear placeholder
                    for row in table.rows:
                        for cell in row.cells:
                            cell.text = cell.text.replace(placeholder, "")
                    
                    # Populate table with data
                    for row_idx, row_data in enumerate(table_data):
                        if row_idx < len(table.rows):
                            for col_idx, cell_value in enumerate(row_data):
                                if col_idx < len(table.rows[row_idx].cells):
                                    table.rows[row_idx].cells[col_idx].text = str(cell_value)
                    
                    return
    
    def _populate_image_field(self, slide, field_name: str, image_path: str):
        """
        Replace a placeholder with an image.
        For now, this is a stub - would require image handling logic.
        """
        # TODO: Implement image insertion
        # This would involve:
        # 1. Finding a shape with the placeholder
        # 2. Getting its position and size
        # 3. Removing the shape
        # 4. Inserting an image at that position
        pass
    
    def get_slide_layouts_info(self, template_path: Path) -> List[Dict[str, Any]]:
        """
        Get information about all slide layouts in a template.
        Returns metadata from slide notes.
        """
        return self.get_template_slide_types(template_path)
    
    def get_template_fields(self, template_path: Path) -> list[str]:
        """
        Get all unique field names from the template.
        Required by base class.
        """
        prs = Presentation(template_path)
        fields = set()
        pattern = r'\{\{([^}]+)\}\}'
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and shape.text:
                    matches = re.findall(pattern, shape.text)
                    fields.update(matches)
                
                if hasattr(shape, "table"):
                    table = shape.table
                    for row in table.rows:
                        for cell in row.cells:
                            matches = re.findall(pattern, cell.text)
                            fields.update(matches)
        
        return sorted(list(fields))
    
    # Keep the original generate method for backward compatibility
    def generate(self, template_path: Path, fields: Dict[str, Any], output_path: Path) -> Path:
        """
        Original generate method - replaces all placeholders in entire template.
        Use generate_from_slides() for advanced slide composition.
        """
        prs = Presentation(template_path)
        
        for slide in prs.slides:
            self._replace_text_in_slide(slide, "{{", "}}")
            # Simple replacement for all fields
            for field_name, field_value in fields.items():
                placeholder = f"{{{{{field_name}}}}}"
                self._replace_text_in_slide(slide, placeholder, str(field_value))
        
        prs.save(output_path)
        return output_path
